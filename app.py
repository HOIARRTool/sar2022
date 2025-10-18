"""
AI-Powered Self-Assessment Report (SAR) Generator
Streamlit app: สร้าง SAR ตามมาตรฐาน HA ฉบับที่ 5 จากไฟล์ที่อัปโหลด (PDF/DOCX/Excel/CSV/PPTX) + ข้อความเสริม
"""

import io
import os
from pathlib import Path
from string import Template
import textwrap

import streamlit as st
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from google import genai

# ====== CONFIGS ======
MAX_ROWS = 300    # จำกัดแถวต่อชีท/ตาราง เพื่อกัน context ล้น
MAX_COLS = 50     # จำกัดคอลัมน์ต่อชีท/ตาราง

# --- Page Configuration ---
st.set_page_config(
    page_title="AI SAR Generator (HA Standard)",
    page_icon="📝",
    layout="wide",
    initial_sidebar_state="expanded"
)
st.caption("Build tag: uploader-xls-pptx-v1")

# --- Hardcoded Knowledge Base & SAR Items ---
SAR_ITEMS = {
    1: "I-1.1ก(1)(2)(3) การชี้นำองค์กรโดยผู้นำระดับสูง",
    2: "I-1.1ข การสื่อสาร สร้างความผูกพันโดยผู้นำ**",
    3: "I-1.1ค(1)(2)(3) การสร้างสิ่งแวดล้อมที่เอื้อต่อการพัฒนาและความสำเร็จขององค์กร**",
    4: "I-1.2ก(1)(2) ระบบกำกับดูแลกิจการ การประเมินผู้นำ/ระบบการนำ",
    5: "I-1.2ก(3) ระบบกำกับดูแลทางคลินิก**",
    6: "I-1.2ข(1)(2)(3),ค(1)(2) การปฏิบัติตามกฎหมาย การทำประโยชน์ให้สังคมและการดำเนินงานอย่างมีจริยธรรม",
    7: "I-2.1ก(1)(2)(3)(4) กระบวนการจัดทำ วางแผนกลยุทธ์และการวิเคราะห์ข้อมูล **",
    8: "I-2.1ข(1)(2)(3) วัตถุประสงค์เชิงกลยุทธ์...",
    9: "I-2.2ก(1)(2)(3)(4) การจัดทำแผนปฏิบัติการ...",
    10: "I-2.2ก(5), ข การกำหนดตัวชี้วัด...",
    11: "I-3.1ก(1) การรับฟัง/เรียนรู้ความต้องการ...",
    # ... (คงเดิมถึง 82)
    82: "III-6 การดูแลต่อเนื่อง**"
}

KNOWLEDGE_BASE = """
- HA Standard 5th Edition: ...
- SPA Manuals (Part I, II, III): ...
- 3P Framework (Purpose-Process-Performance): ...
- SAR Structure: ...
"""

# ===================== Helper: common =====================
def _limit_df(df, max_rows=MAX_ROWS, max_cols=MAX_COLS):
    df = df.copy()
    if df.shape[0] > max_rows:
        df = df.head(max_rows)
    if df.shape[1] > max_cols:
        df = df.iloc[:, :max_cols]
    return df

# ===================== Helper: PDF =====================
def extract_text_from_pdf(file):
    try:
        file.seek(0)
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading PDF file {getattr(file, 'name', '')}: {e}")
        return ""

# ===================== Helper: DOCX =====================
def extract_text_from_docx(file):
    try:
        file.seek(0)
        data = file.read()
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        st.error(f"Error reading DOCX file {getattr(file, 'name', '')}: {e}")
        return ""

# ===================== Helper: CSV =====================
def extract_text_from_csv(file):
    try:
        file.seek(0)
        df = pd.read_csv(file, dtype=str)
        df = _limit_df(df)
        return df.to_csv(index=False)
    except Exception as e:
        st.error(f"Error reading CSV file {getattr(file, 'name', '')}: {e}")
        return ""

# ===================== Helper: Excel =====================
def extract_text_from_excel(file):
    """
    อ่านทุกชีทของ Excel (.xlsx/.xlsm ใช้ openpyxl, .xls ใช้ xlrd)
    """
    try:
        file.seek(0)
        ext = Path(getattr(file, "name", "")).suffix.lower()
        if ext in [".xlsx", ".xlsm"]:
            engine = "openpyxl"
        elif ext == ".xls":
            engine = "xlrd"
        else:
            engine = None  # เผื่อกรณีพิเศษ

        xls = pd.ExcelFile(file, engine=engine) if engine else pd.ExcelFile(file)
        parts = []
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name=sheet_name, dtype=str)
            df = _limit_df(df)
            csv_text = df.to_csv(index=False)
            parts.append(f"[Sheet: {sheet_name}]\n{csv_text}\n")
        return "\n".join(parts)

    except ImportError as ie:
        st.error("ขาดไลบรารีอ่าน Excel:\n• .xlsx/.xlsm ต้องมี openpyxl\n• .xls ต้องมี xlrd")
        st.exception(ie)
        return ""
    except Exception as e:
        st.error(f"Error reading Excel file {getattr(file, 'name', '')}: {e}")
        return ""

# ===================== Helper: PPTX =====================
def _extract_text_from_table(tbl):
    rows = []
    for r in tbl.rows:
        rows.append(",".join(c.text.strip() for c in r.cells))
    return "\n".join(rows)

def _walk_shapes(shapes, out_parts):
    for sh in shapes:
        # text box / placeholder
        if hasattr(sh, "has_text_frame") and sh.has_text_frame:
            txt = sh.text_frame.text or ""
            if txt.strip():
                out_parts.append(txt.strip())
        # table
        if hasattr(sh, "has_table") and sh.has_table:
            out_parts.append(_extract_text_from_table(sh.table))
        # group
        if hasattr(sh, "shapes"):
            _walk_shapes(sh.shapes, out_parts)

def extract_text_from_pptx(file):
    """อ่านข้อความทุกสไลด์ รวม group/table และ notes"""
    try:
        file.seek(0)
        data = file.read()
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"[Slide {i}]"]
            _walk_shapes(slide.shapes, slide_parts)
            # notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = slide.notes_slide.notes_text_frame.text or ""
                if note.strip():
                    slide_parts.append(f"[Notes]\n{note.strip()}")
            parts.append("\n".join(slide_parts))
        return "\n\n".join(parts)
    except Exception as e:
        st.error(f"Error reading PPTX file {getattr(file, 'name', '')}: {e}")
        return ""

# ===================== Aggregate all input =====================
def get_all_input_text(uploaded_files, additional_context):
    full_text = ""
    if not uploaded_files and not additional_context:
        return "", False

    if uploaded_files:
        for file in uploaded_files:
            full_text += f"--- START OF FILE: {file.name} ---\n\n"
            name = (file.name or "").lower()
            mime = (file.type or "").lower()

            try:
                if name.endswith(".pdf") or "pdf" in mime:
                    full_text += extract_text_from_pdf(file)

                elif name.endswith(".docx") or "wordprocessingml" in mime:
                    full_text += extract_text_from_docx(file)

                elif name.endswith((".xlsx", ".xlsm")) or "officedocument.spreadsheetml" in mime:
                    full_text += extract_text_from_excel(file)

                elif name.endswith(".xls") or mime == "application/vnd.ms-excel":
                    full_text += extract_text_from_excel(file)

                elif name.endswith(".csv") or "text/csv" in mime:
                    full_text += extract_text_from_csv(file)

                elif name.endswith(".pptx") or "officedocument.presentationml.presentation" in mime:
                    full_text += extract_text_from_pptx(file)

                else:
                    st.warning(f"Unsupported file type: {file.name}. Skipping.")
            except Exception as e:
                st.error(f"Error reading file {file.name}: {e}")

            full_text += f"\n\n--- END OF FILE: {file.name} ---\n\n"

    if additional_context:
        full_text += f"--- START OF ADDITIONAL USER NOTES ---\n\n{additional_context}\n\n--- END OF ADDITIONAL USER NOTES ---\n\n"

    return full_text, True

# ===================== Gemini call =====================
def generate_sar_section(api_key, sar_item, context_data):
    """Generates a single SAR section using the Gemini API (Google GenAI SDK)."""
    try:
        client = genai.Client(api_key=api_key)

        prompt_tmpl = Template(textwrap.dedent("""\
            คุณคือผู้เชี่ยวชาญด้านการรับรองคุณภาพโรงพยาบาลในประเทศไทย (Hospital Accreditation) ที่มีความสามารถในการเขียนรายงานประเมินตนเอง (Self-Assessment Report - SAR) ตามมาตรฐาน HA ฉบับที่ 5

            **ภารกิจของคุณ:**
            เขียนรายงาน SAR สำหรับหัวข้อต่อไปนี้:
            **${sar_item}**

            **ฐานความรู้ของคุณ:**
            ${knowledge_base}

            **ข้อมูลนำเข้า (จากไฟล์และข้อความที่ผู้ใช้อัปโหลด):**
            ```
            ${context_data}
            ```
            **คำสั่งหลัก (Rule-Based Process):**

            **ขั้นตอนที่ 1: การวิเคราะห์ความเกี่ยวข้อง (Relevance Analysis)**
            เปรียบเทียบเนื้อหาใน "ข้อมูลนำเข้า" กับหัวข้อ **"${sar_item}"** แล้วตัดสินใจว่าพอเพียงหรือไม่

            **ขั้นตอนที่ 2: ดำเนินการตามผลการวิเคราะห์**

            **กรณีที่ 1: ข้อมูลเพียงพอ**
            เขียนรายงาน 4 ส่วน: (i) บริบท, (ii) ประเด็น/แผนการพัฒนา, (iii) ผลที่โดดเด่น, (iv) ผลลัพธ์
            - (i) บริบท: สรุปภาพรวม/นโยบาย/สถานการณ์ที่เกี่ยวข้อง อ้างอิงจากข้อมูลที่ได้รับ
            - (ii) ประเด็น/แผน:
              1) ฉบับเต็ม: ทุกย่อหน้าเริ่มด้วย "เพื่อ..." และเล่าตามลำดับ วัตถุประสงค์→ทำอะไร→ผลเปลี่ยนแปลง→ช่องว่าง
              2) ฉบับสรุป ≤400 ตัวอักษร หัวข้อ "**ii (สรุป ≤400 ตัวอักษร – ทางเลือก):**"
            - (iii) โดดเด่น: 1–2 ประโยคสั้นๆ เน้นใจความสำคัญเท่านั้น
            - (iv) ผลลัพธ์: สรุปเชิงคุณภาพ/เชิงปริมาณ ปิดด้วยหมายเหตุเรื่อง KPI ย้อนหลัง 3–5 ปีให้ผู้ใช้วิเคราะห์เอง

            **กรณีที่ 2: ข้อมูลไม่เพียงพอ**
            ตอบกลับข้อความนี้เท่านั้น:
            "**[AI Analysis]:** จากการตรวจสอบไฟล์และข้อมูลที่ท่านให้มา ไม่พบเนื้อหาที่เกี่ยวข้องโดยตรงกับหัวข้อ **'${sar_item}'** ครับ/ค่ะ กรุณาอัปโหลดเอกสารที่ตรงกับหัวข้อที่เลือก เพื่อให้ AI สามารถสร้างรายงานได้อย่างถูกต้อง"

            **ข้อบังคับเพิ่มเติม**
            - อ้างอิงเฉพาะข้อมูลนำเข้าที่ให้มา
            - ระบุชื่อไฟล์/ชีท/สไลด์/หน้า เมื่ออ้างอิง
        """))

        prompt = prompt_tmpl.substitute(
            sar_item=sar_item,
            knowledge_base=KNOWLEDGE_BASE,
            context_data=context_data
        )

        resp = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )
        return getattr(resp, "text", None) or getattr(resp, "output_text", None)

    except Exception as e:
        st.error(f"An error occurred while calling the Gemini API: {e}")
        st.error("เคล็ดลับ: ตรวจสอบ API Key/สิทธิ์ และชื่อโมเดลว่ารองรับ generateContent")
        return None

# ===================== UI =====================
st.title("📝 Self-Assessment Report (SAR2022) by AI")
st.markdown("เครื่องมือช่วยสร้างรายงานประเมินตนเองตามมาตรฐาน HA ฉบับที่ 5")
st.markdown("---")

# --- API Key ---
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    try:
        GEMINI_API_KEY = st.secrets["GEMINI_API_KEY"]
    except Exception:
        st.error("⚠️ ยังไม่ได้ตั้งค่า API Key!")
        st.info(
            """
**บน Render:** Settings → Environment → เพิ่ม `GEMINI_API_KEY`
**ในเครื่อง:** `.streamlit/secrets.toml` → GEMINI_API_KEY="YOUR_KEY"
            """
        )
        st.stop()

# --- Sidebar ---
with st.sidebar:
    st.header("1. ป้อนข้อมูล")
    st.success("API Key ถูกโหลดเรียบร้อยแล้ว (ENV/secrets)")

    sar_options = [f"{key}. {value}" for key, value in SAR_ITEMS.items()]
    selected_option_str = st.selectbox(
        "เลือกหัวข้อ SAR ที่ต้องการทำ (1-82)",
        options=sar_options,
        index=None,
        placeholder="เลือกหัวข้อ..."
    )

    # เปิดรับชนิดไฟล์ที่ต้องการทั้งหมด (หรือจะเอา type=[] ก็ได้)
    uploaded_files = st.file_uploader(
        "อัปโหลดไฟล์ที่เกี่ยวข้อง (PDF, DOCX, XLSX/XLSM/XLS, CSV, PPTX)",
        type=["pdf", "docx", "xlsx", "xlsm", "xls", "csv", "pptx"],
        accept_multiple_files=True
    )

    additional_context = st.text_area(
        "เขียนข้อมูลหรือบริบทเพิ่มเติม (ถ้ามี)",
        height=150,
        placeholder="เช่น ประเด็นที่ต้องการเน้นเป็นพิเศษ, ข้อมูลที่ไม่มีในเอกสาร..."
    )

    generate_button = st.button("🚀 สร้างรายงาน SAR", use_container_width=True, type="primary")

# --- Output ---
st.header("2. ผลลัพธ์ (AI-Generated SAR)")

if "report_output" not in st.session_state:
    st.session_state.report_output = ""

if generate_button:
    if not selected_option_str:
        st.error("กรุณาเลือกหัวข้อ SAR ที่ต้องการทำ")
    elif not uploaded_files and not additional_context:
        st.error("กรุณาอัปโหลดไฟล์อย่างน้อย 1 ไฟล์ หรือใส่ข้อมูลเพิ่มเติม")
    else:
        with st.spinner("⏳ AI กำลังอ่านไฟล์และเรียบเรียงรายงาน..."):
            context_data, files_ok = get_all_input_text(uploaded_files, additional_context)
            if files_ok:
                generated_report = generate_sar_section(GEMINI_API_KEY, selected_option_str, context_data)
                if generated_report:
                    st.session_state.report_output = generated_report
                else:
                    st.session_state.report_output = ""
                    st.error("ไม่สามารถสร้างรายงานได้ กรุณาตรวจสอบข้อผิดพลาดและลองใหม่")
            else:
                st.error("ไม่สามารถดึงข้อมูลจากไฟล์ที่อัปโหลดได้ กรุณาตรวจสอบไฟล์อีกครั้ง")

if st.session_state.report_output:
    st.markdown(st.session_state.report_output)
else:
    st.info("กรุณาเลือกหัวข้อ, อัปโหลดไฟล์, และกดปุ่ม 'สร้างรายงาน SAR' เพื่อเริ่มต้น")
